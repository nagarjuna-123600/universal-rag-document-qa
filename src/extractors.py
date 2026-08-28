import io
import json
import re
from pathlib import Path
from typing import Iterable

import pandas as pd
from PIL import Image
from langchain_core.documents import Document

from .models import ParsedUnit


# ============================================================
# OCR ENGINE
# ============================================================

# Do NOT initialize RapidOCR when the application starts.
# It will be initialized only when OCR is actually required.
ocr_engine = None


def get_ocr_engine():
    """
    Create the RapidOCR engine only when OCR is needed.
    This prevents Streamlit Cloud from failing during
    application startup if there is an OCR dependency issue.
    """

    global ocr_engine

    if ocr_engine is None:

        try:
            from rapidocr import RapidOCR

            ocr_engine = RapidOCR()

        except Exception as exc:

            raise RuntimeError(
                f"OCR initialization failed: "
                f"{type(exc).__name__}: {exc}"
            ) from exc

    return ocr_engine


# ============================================================
# TEXT CLEANING
# ============================================================

def clean_text(text: str) -> str:

    text = text.replace("\x00", " ")

    text = re.sub(
        r"[ \t]+",
        " ",
        text
    )

    text = re.sub(
        r"\n{3,}",
        "\n\n",
        text
    )

    return text.strip()


# ============================================================
# CREATE PARSED UNIT
# ============================================================

def _doc(
    text: str,
    **metadata
) -> ParsedUnit | None:

    text = clean_text(text)

    if not text:
        return None

    return ParsedUnit(
        text=text,
        metadata=metadata
    )


# ============================================================
# PDF
# ============================================================

def extract_pdf(
    data: bytes,
    file_name: str
) -> list[ParsedUnit]:

    import fitz

    units = []

    with fitz.open(
        stream=data,
        filetype="pdf"
    ) as pdf:

        if pdf.is_encrypted:

            raise ValueError(
                "Encrypted/password-protected PDFs "
                "are not supported."
            )

        for page_no, page in enumerate(
            pdf,
            start=1
        ):

            # ------------------------------------------------
            # First try normal PDF text extraction
            # ------------------------------------------------

            text = page.get_text("text")

            if clean_text(text):

                unit = _doc(
                    text,
                    page=page_no,
                    source_type="pdf"
                )

                if unit:
                    units.append(unit)

            # ------------------------------------------------
            # If no text -> OCR
            # ------------------------------------------------

            else:

                pix = page.get_pixmap(
                    matrix=fitz.Matrix(
                        1.8,
                        1.8
                    ),
                    alpha=False
                )

                img = Image.open(
                    io.BytesIO(
                        pix.tobytes("png")
                    )
                )

                # Initialize OCR only when required
                try:

                    result = get_ocr_engine()(img)

                except Exception as exc:

                    raise RuntimeError(
                        f"PDF OCR failed on page "
                        f"{page_no}: "
                        f"{type(exc).__name__}: {exc}"
                    ) from exc

                ocr_text = ""

                if result and result.txts:

                    ocr_text = "\n".join(
                        result.txts
                    )

                if clean_text(ocr_text):

                    unit = _doc(
                        ocr_text,
                        page=page_no,
                        source_type="pdf_ocr"
                    )

                    if unit:
                        units.append(unit)

    return units


# ============================================================
# DOCX
# ============================================================

def extract_docx(
    data: bytes
) -> list[ParsedUnit]:

    from docx import Document as DocxDocument

    d = DocxDocument(
        io.BytesIO(data)
    )

    units = []

    # --------------------------------------------------------
    # Paragraphs
    # --------------------------------------------------------

    for i, paragraph in enumerate(
        d.paragraphs,
        start=1
    ):

        if paragraph.text.strip():

            unit = _doc(
                paragraph.text,
                paragraph=i,
                source_type="docx"
            )

            if unit:
                units.append(unit)

    # --------------------------------------------------------
    # Tables
    # --------------------------------------------------------

    for table_index, table in enumerate(
        d.tables,
        start=1
    ):

        rows = [
            " | ".join(
                cell.text.strip()
                for cell in row.cells
            )
            for row in table.rows
        ]

        unit = _doc(
            "\n".join(rows),
            section=f"table {table_index}",
            source_type="docx_table"
        )

        if unit:
            units.append(unit)

    return units


# ============================================================
# XLSX
# ============================================================

def extract_xlsx(
    data: bytes,
    file_name: str
) -> list[ParsedUnit]:

    book = pd.ExcelFile(
        io.BytesIO(data)
    )

    units = []

    for sheet in book.sheet_names:

        df = pd.read_excel(
            book,
            sheet_name=sheet,
            dtype=str
        ).fillna("")

        text = (
            f"Sheet: {sheet}\n"
            + df.to_csv(index=False)
        )

        unit = _doc(
            text,
            sheet=sheet,
            source_type="xlsx"
        )

        if unit:
            units.append(unit)

    return units


# ============================================================
# CSV
# ============================================================

def extract_csv(
    data: bytes
) -> list[ParsedUnit]:

    try:

        df = pd.read_csv(
            io.BytesIO(data),
            dtype=str
        ).fillna("")

    except UnicodeDecodeError:

        df = pd.read_csv(
            io.BytesIO(data),
            dtype=str,
            encoding="latin-1"
        ).fillna("")

    text = df.to_csv(
        index=False
    )

    unit = _doc(
        text,
        section="CSV table",
        source_type="csv"
    )

    return [unit] if unit else []


# ============================================================
# PPTX
# ============================================================

def extract_pptx(
    data: bytes
) -> list[ParsedUnit]:

    from pptx import Presentation

    prs = Presentation(
        io.BytesIO(data)
    )

    units = []

    for slide_no, slide in enumerate(
        prs.slides,
        start=1
    ):

        texts = []

        for shape in slide.shapes:

            if (
                hasattr(shape, "text")
                and shape.text.strip()
            ):

                texts.append(
                    shape.text
                )

        if texts:

            unit = _doc(
                "\n".join(texts),
                slide=slide_no,
                source_type="pptx"
            )

            if unit:
                units.append(unit)

    return units


# ============================================================
# JSON
# ============================================================

def extract_json(
    data: bytes
) -> list[ParsedUnit]:

    obj = json.loads(
        data.decode("utf-8-sig")
    )

    text = json.dumps(
        obj,
        ensure_ascii=False,
        indent=2
    )

    unit = _doc(
        text,
        section="JSON document",
        source_type="json"
    )

    return [unit] if unit else []


# ============================================================
# XML
# ============================================================

def extract_xml(
    data: bytes
) -> list[ParsedUnit]:

    from lxml import etree

    root = etree.fromstring(data)

    text = "\n".join(
        t.strip()
        for t in root.itertext()
        if t.strip()
    )

    unit = _doc(
        text,
        section="XML document",
        source_type="xml"
    )

    return [unit] if unit else []


# ============================================================
# TXT
# ============================================================

def extract_text(
    data: bytes,
    file_name: str
) -> list[ParsedUnit]:

    text = data.decode(
        "utf-8-sig",
        errors="strict"
    )

    unit = _doc(
        text,
        section="Text document",
        source_type="txt"
    )

    return [unit] if unit else []


# ============================================================
# IMAGE / OCR
# ============================================================

def extract_image(
    data: bytes
) -> list[ParsedUnit]:

    # --------------------------------------------------------
    # Validate image
    # --------------------------------------------------------

    try:

        img = Image.open(
            io.BytesIO(data)
        )

        img.verify()

    except Exception as exc:

        raise ValueError(
            f"Invalid or corrupted image: {exc}"
        )

    # --------------------------------------------------------
    # Re-open after verify()
    # --------------------------------------------------------

    img = Image.open(
        io.BytesIO(data)
    )

    # --------------------------------------------------------
    # Run OCR
    # --------------------------------------------------------

    try:

        result = get_ocr_engine()(img)

    except Exception as exc:

        raise RuntimeError(
            f"Image OCR failed: "
            f"{type(exc).__name__}: {exc}"
        ) from exc

    # --------------------------------------------------------
    # Extract OCR text
    # --------------------------------------------------------

    text = ""

    if result and result.txts:

        text = "\n".join(
            result.txts
        )

    text = clean_text(text)

    # --------------------------------------------------------
    # No text found
    # --------------------------------------------------------

    if not text:

        raise ValueError(
            "No readable text was found in the image."
        )

    # --------------------------------------------------------
    # Create parsed unit
    # --------------------------------------------------------

    unit = ParsedUnit(
        text=text,
        metadata={
            "section": "OCR image",
            "source_type": "image_ocr"
        }
    )

    return [unit]


# ============================================================
# MAIN DOCUMENT EXTRACTOR
# ============================================================

def extract_document(
    name: str,
    data: bytes
) -> list[ParsedUnit]:

    ext = Path(
        name
    ).suffix.lower()

    # --------------------------------------------------------
    # PDF
    # --------------------------------------------------------

    if ext == ".pdf":

        return extract_pdf(
            data,
            name
        )

    # --------------------------------------------------------
    # DOCX
    # --------------------------------------------------------

    if ext == ".docx":

        return extract_docx(
            data
        )

    # --------------------------------------------------------
    # CSV
    # --------------------------------------------------------

    if ext == ".csv":

        return extract_csv(
            data
        )

    # --------------------------------------------------------
    # XLSX
    # --------------------------------------------------------

    if ext == ".xlsx":

        return extract_xlsx(
            data,
            name
        )

    # --------------------------------------------------------
    # PPTX
    # --------------------------------------------------------

    if ext == ".pptx":

        return extract_pptx(
            data
        )

    # --------------------------------------------------------
    # JSON
    # --------------------------------------------------------

    if ext == ".json":

        return extract_json(
            data
        )

    # --------------------------------------------------------
    # XML
    # --------------------------------------------------------

    if ext == ".xml":

        return extract_xml(
            data
        )

    # --------------------------------------------------------
    # TXT
    # --------------------------------------------------------

    if ext == ".txt":

        return extract_text(
            data,
            name
        )

    # --------------------------------------------------------
    # IMAGE
    # --------------------------------------------------------

    if ext in {
        ".jpg",
        ".jpeg",
        ".png",
        ".webp",
        ".tif",
        ".tiff"
    }:

        return extract_image(
            data
        )

    # --------------------------------------------------------
    # Legacy formats
    # --------------------------------------------------------

    if ext in {
        ".doc",
        ".xls",
        ".ppt"
    }:

        raise ValueError(
            f"Legacy {ext} files require "
            "LibreOffice/antiword conversion "
            "and are not enabled in this safe baseline."
        )

    # --------------------------------------------------------
    # Unsupported
    # --------------------------------------------------------

    raise ValueError(
        f"No extractor configured for {ext}."
    )


# ============================================================
# CONVERT PARSED UNITS → LANGCHAIN DOCUMENTS
# ============================================================

def units_to_documents(
    units: Iterable[ParsedUnit],
    file_name: str,
    file_id: str
) -> list[Document]:

    docs = []

    for idx, unit in enumerate(units):

        # Safety check
        if unit is None:
            continue

        metadata = {
            **unit.metadata,
            "file_name": file_name,
            "file_id": file_id,
            "unit_id": idx,
        }

        docs.append(
            Document(
                page_content=unit.text,
                metadata=metadata
            )
        )

    return docs
